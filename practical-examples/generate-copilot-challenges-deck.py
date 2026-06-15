from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path(__file__).with_name("microsoft-copilot-challenges-agent-comparison.pptx")

NAVY = RGBColor(18, 32, 49)
TEAL = RGBColor(0, 128, 136)
MINT = RGBColor(199, 232, 221)
LIME = RGBColor(214, 230, 158)
CORAL = RGBColor(242, 151, 126)
AMBER = RGBColor(246, 190, 88)
INK = RGBColor(31, 42, 55)
MUTED = RGBColor(92, 105, 120)
PAPER = RGBColor(248, 250, 247)
WHITE = RGBColor(255, 255, 255)
LINK = RGBColor(0, 97, 161)


CHALLENGES = [
    {
        "num": "01",
        "theme": "Qualite & securite",
        "title": "Renforcer le traitement des incidents",
        "answer": "Security Copilot est l'agent cle pour les incidents securite. L'observability agent aide plutot les SRE a analyser les signaux applicatifs et infra.",
        "solutions": [
            {
                "name": "Microsoft Security Copilot",
                "demo": "https://adoption.microsoft.com/en-us/microsoft-security-copilot/",
                "video": "https://learn-video.azurefd.net/vod/player?id=257c1a95-04fb-422e-add2-51e91388ae5f",
            },
            {
                "name": "Azure Copilot observability agent",
                "demo": "https://learn.microsoft.com/en-us/azure/azure-monitor/aiops/observability-agent-overview",
            },
        ],
        "github": {
            "name": "GitHub Copilot agent",
            "demo": "https://docs.github.com/en/copilot/how-tos/use-copilot-agents/cloud-agent",
            "note": "Alternative pour automatiser runbooks, KQL et integrations ITSM depuis le code.",
        },
        "comparison": {
            "other": "Pret a l'emploi pour triage secu, resumes d'incident, KQL et guidage SOC/SRE.",
            "github": "Plus d'effort: connecter logs, runbooks, tickets et garde-fous avant d'obtenir un agent operationnel.",
        },
        "diagram": ["Alertes", "Analyse", "Runbook", "Action"],
        "value": ["Resume l'incident", "Propose la requete KQL", "Guide la remediation"],
        "sources": ["learn.microsoft.com/copilot/security/microsoft-security-copilot", "learn.microsoft.com/azure/sentinel/overview", "learn.microsoft.com/azure/azure-monitor/fundamentals/overview"],
    },
    {
        "num": "02",
        "theme": "Qualite & securite",
        "title": "Securiser et fiabiliser le traitement des donnees",
        "answer": "Passer d'une chaine de donnees dispersee a une plateforme gouvernee: classification, controles, lineage, supervision et usages IA bornes.",
        "solutions": [
            {
                "name": "Copilot in Microsoft Purview",
                "demo": "https://learn.microsoft.com/en-us/purview/ai-microsoft-purview",
            },
            {
                "name": "Copilot in Fabric",
                "demo": "https://learn.microsoft.com/en-us/fabric/fundamentals/copilot-fabric-overview",
            },
        ],
        "github": {
            "name": "GitHub Copilot agent",
            "demo": "https://docs.github.com/en/copilot/how-tos/use-copilot-agents/cloud-agent",
            "note": "Alternative pour coder controles, pipelines et rapports autour du patrimoine data.",
        },
        "comparison": {
            "other": "Purview/Fabric donnent deja classification, gouvernance, lineage et aide a l'analyse.",
            "github": "A construire: connecteurs, scripts de controle, tests et documentation des regles data.",
        },
        "diagram": ["Donnees", "Classement", "Controle", "Usage"],
        "value": ["Identifie les donnees sensibles", "Explique les risques", "Aide a documenter les controles"],
        "sources": ["learn.microsoft.com/purview/purview", "learn.microsoft.com/fabric/fundamentals/microsoft-fabric-overview", "learn.microsoft.com/azure/defender-for-cloud/defender-for-cloud-introduction"],
    },
    {
        "num": "03",
        "theme": "Qualite & securite",
        "title": "Piloter les plans de remediation SSI par les risques",
        "answer": "Prioriser les actions SSI par criticite business, chemin d'attaque, exposition cloud/code et capacite de remediation.",
        "solutions": [
            {
                "name": "Microsoft Security Copilot",
                "demo": "https://adoption.microsoft.com/en-us/microsoft-security-copilot/",
                "video": "https://learn-video.azurefd.net/vod/player?id=257c1a95-04fb-422e-add2-51e91388ae5f",
            },
            {
                "name": "Defender for Cloud embedded Copilot",
                "demo": "https://learn.microsoft.com/en-us/copilot/security/experiences-security-copilot#embedded-experiences",
            },
        ],
        "github": {
            "name": "GitHub Copilot agent",
            "demo": "https://docs.github.com/en/copilot/how-tos/use-copilot-agents/cloud-agent",
            "note": "Alternative pour generer correctifs IaC/code et pull requests de remediation.",
        },
        "comparison": {
            "other": "Security Copilot/Defender savent analyser recommandations, attaque paths et priorites risque.",
            "github": "Utile pour transformer les actions en code, mais demande mapping repo, policies et validations.",
        },
        "diagram": ["Risque", "Priorite", "Owner", "Suivi"],
        "value": ["Explique le chemin d'attaque", "Priorise par impact", "Prepare les actions owners"],
        "sources": ["learn.microsoft.com/azure/defender-for-cloud/defender-for-cloud-introduction", "learn.microsoft.com/copilot/security/microsoft-security-copilot"],
    },
    {
        "num": "04",
        "theme": "Relation metier / IT",
        "title": "Alleger le pilotage des projets",
        "answer": "Automatiser la synthese, les comptes rendus, la preparation des supports et le suivi des decisions sans sortir du flux Teams/Office.",
        "solutions": [
            {
                "name": "Microsoft 365 Copilot",
                "demo": "https://adoption.microsoft.com/en-us/copilot/",
            },
            {
                "name": "Copilot Studio agent",
                "demo": "https://learn.microsoft.com/en-us/microsoft-copilot-studio/fundamentals-what-is-copilot-studio",
            },
            {
                "name": "Azure AI Foundry Agent Service",
                "demo": "https://learn.microsoft.com/en-us/azure/ai-foundry/agents/overview",
            },
        ],
        "github": {
            "name": "GitHub Copilot agent",
            "demo": "https://docs.github.com/en/copilot/how-tos/use-copilot-agents/cloud-agent",
            "note": "Alternative pro-code pour construire/connecter un agent projet sur mesure.",
        },
        "comparison": {
            "other": "Low code: M365/Copilot Studio pour Teams, actions et workflows. Pro code: Foundry pour agent custom.",
            "github": "GitHub Copilot aide a coder l'agent Foundry, connecteurs et tests; plus flexible, plus d'engineering.",
        },
        "diagram": ["Reunion", "Synthese", "Actions", "Reporting"],
        "value": ["Produit le CR", "Suit les decisions", "Prepare le support projet"],
        "sources": ["learn.microsoft.com/microsoft-365/copilot/microsoft-365-copilot-overview", "learn.microsoft.com/azure/ai-foundry/agents/overview"],
    },
    {
        "num": "05",
        "theme": "Efficience & patrimoine",
        "title": "Maitriser les patrimoines applicatifs",
        "answer": "Construire un inventaire vivant des applications, dependances, couts, risques et trajectoires de transformation.",
        "solutions": [
            {
                "name": "Azure Copilot migration agent",
                "demo": "https://learn.microsoft.com/en-us/azure/migrate/azure-copilot-migration-agent?view=migrate",
            },
            {
                "name": "Azure AI Foundry Agent Service",
                "demo": "https://learn.microsoft.com/en-us/azure/ai-foundry/agents/overview",
            },
        ],
        "github": {
            "name": "GitHub Copilot agent",
            "demo": "https://docs.github.com/en/copilot/how-tos/use-copilot-agents/cloud-agent",
            "note": "Alternative pour generer fiches, scripts Resource Graph et backlog de rationalisation.",
        },
        "comparison": {
            "other": "Azure Migrate lit deja l'inventaire. Foundry peut porter un agent connecte a Azure Resource Graph.",
            "github": "A construire: requetes Resource Graph, schema applicatif, scoring et generation de backlog.",
        },
        "diagram": ["Inventaire", "Dependances", "Trajectoire", "Backlog"],
        "value": ["Interroge l'inventaire", "Resume les dependances", "Cree une fiche applicative"],
        "sources": ["learn.microsoft.com/azure/migrate/migrate-services-overview", "learn.microsoft.com/azure/cloud-adoption-framework/", "learn.microsoft.com/azure/ai-foundry/agents/overview"],
    },
    {
        "num": "06",
        "theme": "Socles technologiques",
        "title": "Accelerer la modernisation des patrimoines coeurs",
        "answer": "Industrialiser la trajectoire migrate/modernize: evaluation, refactoring assiste, landing zones et pipeline de livraison securise.",
        "solutions": [
            {
                "name": "Azure Copilot migration agent",
                "demo": "https://learn.microsoft.com/en-us/azure/migrate/azure-copilot-migration-agent?view=migrate",
            },
            {
                "name": "GitHub Copilot coding agent",
                "demo": "https://docs.github.com/en/copilot/how-tos/use-copilot-agents/cloud-agent",
            },
        ],
        "github": {
            "name": "GitHub Copilot agent",
            "demo": "https://docs.github.com/en/copilot/how-tos/use-copilot-agents/cloud-agent",
            "note": "Alternative directe pour refactorer code, IaC, tests et pipelines.",
        },
        "comparison": {
            "other": "Azure Migrate planifie la trajectoire et cible Azure; l'execution code reste a industrialiser.",
            "github": "Fort pour pro-code: PR, refactoring, IaC et tests; demande contexte repo et validation humaine.",
        },
        "diagram": ["Evaluer", "Planifier", "Refactorer", "Deployer"],
        "value": ["Compare les options", "Genere code/IaC/tests", "Prepare la migration progressive"],
        "sources": ["learn.microsoft.com/azure/migrate/migrate-services-overview", "docs.github.com/copilot", "learn.microsoft.com/azure/cloud-adoption-framework/"],
    },
    {
        "num": "07",
        "theme": "Transverse",
        "title": "Simplifier la gestion budgetaire",
        "answer": "Mettre une boucle FinOps continue: allocation, prevision, alertes, optimisation et arbitrage partage avec les equipes.",
        "solutions": [
            {
                "name": "Microsoft Copilot for Azure",
                "demo": "https://learn.microsoft.com/en-us/azure/copilot/overview",
            },
            {
                "name": "FinOps toolkit reporting agent",
                "demo": "https://learn.microsoft.com/en-us/cloud-computing/finops/toolkit/power-bi/reports",
            },
        ],
        "github": {
            "name": "GitHub Copilot agent",
            "demo": "https://docs.github.com/en/copilot/how-tos/use-copilot-agents/cloud-agent",
            "note": "Alternative pour automatiser exports, controles FinOps et rapports custom.",
        },
        "comparison": {
            "other": "Copilot for Azure et FinOps Toolkit donnent vues, anomalies, reporting et optimisation standard.",
            "github": "A construire: jobs, dashboards custom, tests de tags et PR de correction sur IaC.",
        },
        "diagram": ["Couts", "Anomalies", "Budget", "Action"],
        "value": ["Explique les ecarts", "Prepare le reporting", "Suggere les optimisations"],
        "sources": ["learn.microsoft.com/azure/cost-management-billing/costs/quick-acm-cost-analysis", "learn.microsoft.com/cloud-computing/finops/toolkit/finops-toolkit-overview"],
    },
]


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def add_text(slide, text, x, y, w, h, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_pill(slide, text, x, y, w, color):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.38))
    set_fill(pill, color)
    pill.text_frame.text = text
    pill.text_frame.margin_left = Inches(0.04)
    pill.text_frame.margin_right = Inches(0.04)
    pill.text_frame.margin_top = Inches(0.02)
    pill.text_frame.margin_bottom = Inches(0.02)
    pill.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    font_size = 9 if len(text) <= 22 else 7.5
    for p in pill.text_frame.paragraphs:
        for r in p.runs:
            r.font.name = "Aptos"
            r.font.size = Pt(font_size)
            r.font.bold = True
            r.font.color.rgb = NAVY
    return pill


def add_link(slide, text, url, x, y, w, h, size=8.5):
    box = add_text(slide, text, x, y, w, h, size, LINK, True)
    run = box.text_frame.paragraphs[0].runs[0]
    run.font.underline = True
    run.hyperlink.address = url
    return box


def add_card(slide, title, body, x, y, w, h, fill, title_color=NAVY):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(card, fill)
    card.line.color.rgb = MINT
    add_text(slide, title, x + 0.13, y + 0.11, w - 0.26, 0.18, 9.4, title_color, True)
    add_text(slide, body, x + 0.13, y + 0.35, w - 0.26, h - 0.43, 7.2, INK)
    return card


def add_flow(slide, labels, x, y, w, color):
    step_w = w / len(labels)
    for i, label in enumerate(labels):
        sx = x + i * step_w
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(sx), Inches(y), Inches(step_w - 0.12), Inches(0.62))
        set_fill(shape, color if i % 2 == 0 else MINT)
        shape.text_frame.text = label
        shape.text_frame.margin_left = Inches(0.05)
        shape.text_frame.margin_right = Inches(0.05)
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = "Aptos"
                r.font.size = Pt(9)
                r.font.bold = True
                r.font.color.rgb = NAVY
        if i < len(labels) - 1:
            slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(sx + step_w - 0.17),
                Inches(y + 0.22),
                Inches(0.14),
                Inches(0.16),
            ).fill.solid()
            arrow = slide.shapes[-1]
            arrow.fill.fore_color.rgb = TEAL
            arrow.line.color.rgb = TEAL


def add_challenge_slide(prs, item, index):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.22))
    set_fill(band, TEAL if index % 2 else NAVY)

    add_text(slide, f"Challenge #{item['num']}", 0.55, 0.48, 1.35, 0.3, 12, TEAL, True)
    add_text(slide, item["theme"].upper(), 7.05, 0.5, 2.4, 0.25, 8.5, MUTED, True, PP_ALIGN.RIGHT)
    add_text(slide, item["title"], 0.55, 0.78, 6.8, 0.72, 20, NAVY, True)
    add_text(slide, item["answer"], 0.58, 1.56, 6.75, 0.86, 9.4, INK)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.25), Inches(0.82), Inches(0.92), Inches(0.92))
    set_fill(circle, LIME if index in (2, 6) else MINT)
    add_text(slide, item["num"], 8.43, 1.06, 0.55, 0.24, 15, NAVY, True, PP_ALIGN.CENTER)

    add_text(slide, "Agents specialises", 0.58, 2.43, 2.3, 0.25, 12.5, NAVY, True)
    for i, solution in enumerate(item["solutions"][:3]):
        color = [MINT, LIME, CORAL][i]
        y = 2.78 + i * 0.42
        pill = add_pill(slide, solution["name"], 0.58, y, 3.1, color)
        pill.click_action.hyperlink.address = solution["demo"]
        add_link(slide, "Demo", solution["demo"], 3.82, y + 0.09, 0.55, 0.16, 7.4)
        if "video" in solution:
            add_link(slide, "Video", solution["video"], 4.3, y + 0.09, 0.55, 0.16, 7.4)

    add_text(slide, "Alternative GitHub", 5.05, 2.43, 2.3, 0.25, 12.5, NAVY, True)
    github = item["github"]
    github_pill = add_pill(slide, github["name"], 5.05, 2.78, 3.15, MINT)
    github_pill.click_action.hyperlink.address = github["demo"]
    add_link(slide, "Demo", github["demo"], 8.34, 2.87, 0.55, 0.16, 7.4)
    add_text(slide, github["note"], 5.08, 3.28, 4.05, 0.42, 8.4, MUTED)

    add_card(slide, "Disponible dans les agents", item["comparison"]["other"], 0.58, 3.95, 4.25, 1.27, WHITE)
    add_card(slide, "A faire avec GitHub Copilot", item["comparison"]["github"], 5.05, 3.95, 4.35, 1.27, WHITE)


def add_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_text(slide, "Microsoft Copilot & Azure", 0.65, 0.62, 6.8, 0.45, 20, MINT, True)
    add_text(slide, "Agents specialises vs GitHub Copilot", 0.65, 1.35, 8.4, 0.85, 30, WHITE, True)
    add_text(slide, "Comparer ce qui existe deja et ce qu'il faut construire ou adapter", 0.68, 2.35, 8.2, 0.4, 14, WHITE)
    add_flow(slide, ["Incidents", "Data", "Risques", "Projets", "Apps", "Modernize", "Budget"], 0.72, 3.55, 8.6, LIME)
    add_text(slide, "Slides synthetiques, orientees agents et actions.", 0.7, 5.12, 8.6, 0.2, 9, MINT)


def add_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_text(slide, "Lecture transversale", 0.58, 0.45, 5.2, 0.5, 28, NAVY, True)
    rows = [
        ("Run", "Security Copilot pour incidents secu; observability agent pour SRE"),
        ("Secure", "Purview/Fabric/Defender si disponible; GitHub pour controles code"),
        ("Build", "Foundry pour agent pro-code; Copilot Studio pour low-code"),
        ("Move", "Azure Migrate pour inventaire; GitHub pour backlog/code"),
        ("Optimize", "Copilot for Azure/FinOps pour analyse; GitHub pour automatiser"),
    ]
    for i, (label, text) in enumerate(rows):
        y = 1.25 + i * 0.72
        tag = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(y), Inches(1.22), Inches(0.42))
        set_fill(tag, [MINT, LIME, CORAL, AMBER, MINT][i])
        add_text(slide, label, 0.86, y + 0.1, 0.95, 0.16, 9.5, NAVY, True, PP_ALIGN.CENTER)
        add_text(slide, text, 2.18, y + 0.07, 6.8, 0.25, 14, INK)
    add_text(slide, "Principe", 0.72, 5.0, 2.0, 0.2, 12, NAVY, True)
    add_text(slide, "Choisir l'agent specialise quand il sait deja faire; utiliser GitHub Copilot quand il faut coder, connecter ou industrialiser.", 2.55, 4.92, 6.8, 0.42, 10, MUTED)


def add_sources(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_text(slide, "Sources et diagrammes Microsoft a consulter", 0.55, 0.42, 8.8, 0.42, 24, NAVY, True)
    sources = []
    for item in CHALLENGES:
        for source in item["sources"]:
            if source not in sources:
                sources.append(source)
    for i, source in enumerate(sources[:14]):
        add_text(slide, source, 0.72, 1.0 + i * 0.28, 8.8, 0.16, 7.7, MUTED)
    add_text(slide, "Les schemas inclus sont des versions simplifiees inspirees des architectures de reference, pas des captures de documentation.", 0.72, 5.02, 8.8, 0.38, 8.5, TEAL, True)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    add_title(prs)
    for index, challenge in enumerate(CHALLENGES, start=1):
        add_challenge_slide(prs, challenge, index)
    add_summary(prs)
    add_sources(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()